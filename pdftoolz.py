from flask import Flask, render_template_string, request, send_file, jsonify
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
import pdfplumber
import pandas as pd
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import os
import io
import zipfile
import uuid
import time

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['DOWNLOAD_FOLDER'] = 'downloads'

os.makedirs('uploads', exist_ok=True)
os.makedirs('downloads', exist_ok=True)

HTML = """
<!DOCTYPE html>
<html lang="en" class="scroll-smooth">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0"/>
  <title>PDFToolz - Free Online PDF Tools</title>
  <script src="https://cdn.tailwindcss.com"></script>
  <script src="https://unpkg.com/@phosphor-icons/web"></script>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
  <script>
    tailwind.config = { darkMode: 'class' }
  </script>
  <style>
    body { font-family: 'Inter', sans-serif; }
    .drop-zone { transition: all 0.3s; }
    .drop-zone.drag-active { background: #f0f9ff; border-color: #0ea5e9; transform: scale(1.02); }
    .dark .drop-zone.drag-active { background: #1e293b; }
    @keyframes spin { to { transform: rotate(360deg); } }
    .spinner { animation: spin 1s linear infinite; }
  </style>
</head>
<body class="bg-gray-50 dark:bg-slate-900 text-gray-800 dark:text-gray-100 min-h-screen flex flex-col">

  <!-- Navbar -->
  <nav class="bg-white dark:bg-slate-800 shadow-sm border-b border-gray-200 dark:border-slate-700 sticky top-0 z-50">
    <div class="max-w-7xl mx-auto px-4 sm:px-6 lg:px-8 py-4 flex justify-between items-center">
      <div class="flex items-center gap-3">
        <div class="bg-blue-600 text-white p-2 rounded-lg">
          <i class="ph ph-files text-2xl"></i>
        </div>
        <h1 class="text-2xl font-bold">PDF<span class="text-blue-600">Toolz</span></h1>
      </div>
      <button onclick="document.documentElement.classList.toggle('dark')" class="p-2 rounded-lg hover:bg-gray-100 dark:hover:bg-slate-700">
        <i class="ph ph-moon text-xl"></i>
      </button>
    </div>
  </nav>

  <!-- Hero -->
  <div class="bg-gradient-to-b from-blue-50 to-white dark:from-slate-900 dark:to-slate-900 py-20 text-center">
    <h1 class="text-5xl font-bold mb-4">All-in-One Free PDF Tools</h1>
    <p class="text-xl text-gray-600 dark:text-gray-400 max-w-3xl mx-auto">
      Merge, split, compress, convert PDF to Word, Excel, PowerPoint — 100% free, no signup, no limits.
    </p>
  </div>

  <!-- Tools Grid -->
  <div class="max-w-7xl mx-auto px-6 py-12 grid grid-cols-1 sm:grid-cols-2 lg:grid-cols-3 xl:grid-cols-4 gap-8">
    {% for tool in tools %}
    <div onclick="openTool('{{ tool.id }}')" class="bg-white dark:bg-slate-800 rounded-2xl p-8 shadow-lg hover:shadow-2xl cursor-pointer transform hover:-translate-y-2 transition-all border border-gray-200 dark:border-slate-700">
      <div class="w-16 h-16 {{ tool.bg }} rounded-xl flex items-center justify-center mb-6">
        <i class="ph-fill {{ tool.icon }} text-4xl {{ tool.color }}"></i>
      </div>
      <h3 class="text-xl font-bold mb-2">{{ tool.name }}</h3>
      <p class="text-gray-600 dark:text-gray-400 text-sm">{{ tool.desc }}</p>
    </div>
    {% endfor %}
  </div>

  <!-- Footer -->
  <footer class="mt-auto py-8 text-center text-gray-500 dark:text-gray-400 text-sm">
    <p>PDFToolz - All processing happens securely. Your files are automatically deleted after processing.</p>
  </footer>

  <!-- Tool Modal -->
  <div id="tool-modal" class="fixed inset-0 bg-black/50 z-50 hidden items-center justify-center p-4 overflow-y-auto">
    <div class="bg-white dark:bg-slate-800 rounded-3xl shadow-2xl max-w-2xl w-full p-8 relative my-8">
      <button onclick="closeModal()" class="absolute top-4 right-4 text-gray-500 hover:text-gray-800 dark:hover:text-gray-200">
        <i class="ph ph-x text-2xl"></i>
      </button>
      <h2 id="modal-title" class="text-3xl font-bold mb-6"></h2>

      <!-- Upload Area -->
      <div id="drop-zone" class="drop-zone border-4 border-dashed border-gray-300 dark:border-slate-600 rounded-2xl p-12 text-center cursor-pointer">
        <input type="file" id="file-input" class="hidden" accept=".pdf"/>
        <i class="ph ph-cloud-arrow-up text-6xl text-gray-400 mb-4 block"></i>
        <p class="text-xl font-semibold">Drop PDF here or click to upload</p>
        <p class="text-sm text-gray-500 mt-2">Max file size: 500MB</p>
      </div>

      <!-- File Info -->
      <div id="file-info" class="hidden mt-6 p-4 bg-gray-100 dark:bg-slate-700 rounded-xl">
        <div id="file-list"></div>
        <button onclick="removeFile()" class="text-red-600 text-sm mt-2 hover:underline">Remove all files</button>
      </div>

      <!-- Options Form -->
      <form id="options-form" class="mt-6 space-y-4"></form>

      <!-- Buttons -->
      <div class="mt-8 flex gap-4 justify-end">
        <button type="button" onclick="closeModal()" class="px-6 py-3 bg-gray-300 dark:bg-slate-700 rounded-xl font-medium hover:bg-gray-400 dark:hover:bg-slate-600 transition">Cancel</button>
        <button type="button" id="process-btn" class="px-8 py-3 bg-blue-600 text-white rounded-xl font-bold flex items-center gap-2 hover:bg-blue-700 transition">
          <span id="btn-text">Process PDF</span>
          <i id="spinner" class="ph ph-spinner spinner hidden"></i>
        </button>
      </div>

      <!-- Download -->
      <a id="download-link" class="hidden mt-6 inline-flex items-center gap-3 px-8 py-4 bg-green-600 text-white rounded-xl font-bold text-lg hover:bg-green-700 transition">
        <i class="ph ph-download"></i> Download File
      </a>
      
      <!-- Success Message -->
      <div id="success-msg" class="hidden mt-4 p-4 bg-green-100 dark:bg-green-900 text-green-800 dark:text-green-200 rounded-xl text-center">
        <i class="ph ph-check-circle text-2xl mb-2"></i>
        <p>Processing complete! Click the download button above.</p>
      </div>
    </div>
  </div>

  <script>
    const tools = {{ tools|tojson }};
    let current_tool = null;
    let uploadedFiles = [];

    function openTool(id) {
      current_tool = tools.find(t => t.id === id);
      document.getElementById('modal-title').textContent = current_tool.name;
      document.getElementById('tool-modal').classList.remove('hidden');
      document.getElementById('tool-modal').classList.add('flex');
      
      const optionsForm = document.getElementById('options-form');
      optionsForm.innerHTML = current_tool.options || '';
      if (current_tool.options) {
        optionsForm.classList.remove('hidden');
      } else {
        optionsForm.classList.add('hidden');
      }
      
      const fileInput = document.getElementById('file-input');
      fileInput.multiple = current_tool.multiple || false;
      
      resetUI();
    }

    function closeModal() {
      document.getElementById('tool-modal').classList.add('hidden');
      document.getElementById('tool-modal').classList.remove('flex');
      resetUI();
    }

    function resetUI() {
      uploadedFiles = [];
      document.getElementById('file-input').value = '';
      document.getElementById('file-info').classList.add('hidden');
      document.getElementById('file-list').innerHTML = '';
      document.getElementById('download-link').classList.add('hidden');
      document.getElementById('success-msg').classList.add('hidden');
      document.getElementById('btn-text').textContent = 'Process PDF';
      document.getElementById('spinner').classList.add('hidden');
      document.getElementById('process-btn').disabled = false;
    }

    function removeFile() {
      uploadedFiles = [];
      document.getElementById('file-input').value = '';
      document.getElementById('file-info').classList.add('hidden');
      document.getElementById('file-list').innerHTML = '';
    }

    // Drag & Drop
    const dropZone = document.getElementById('drop-zone');
    
    ['dragenter', 'dragover'].forEach(eventName => {
      dropZone.addEventListener(eventName, (e) => {
        e.preventDefault();
        e.stopPropagation();
        dropZone.classList.add('drag-active');
      });
    });
    
    ['dragleave', 'drop'].forEach(eventName => {
      dropZone.addEventListener(eventName, (e) => {
        e.preventDefault();
        e.stopPropagation();
        dropZone.classList.remove('drag-active');
      });
    });
    
    dropZone.addEventListener('drop', e => {
      const files = e.dataTransfer.files;
      handleFiles(files);
    });
    
    dropZone.addEventListener('click', () => document.getElementById('file-input').click());
    document.getElementById('file-input').addEventListener('change', e => handleFiles(e.target.files));

    function handleFiles(files) {
      if (files.length === 0) return;
      
      const validFiles = [];
      for (let file of files) {
        if (file.name.toLowerCase().endsWith('.pdf')) {
          validFiles.push(file);
        }
      }
      
      if (validFiles.length === 0) {
        alert('Please upload PDF files only');
        return;
      }
      
      uploadedFiles = validFiles;
      
      const fileList = document.getElementById('file-list');
      fileList.innerHTML = '';
      
      validFiles.forEach((file, index) => {
        const div = document.createElement('div');
        div.className = 'flex items-center gap-2 py-1';
        div.innerHTML = `
          <i class="ph ph-file-pdf text-red-500"></i>
          <span class="font-medium">${file.name}</span>
          <span class="text-gray-500 text-sm">(${(file.size/1024/1024).toFixed(2)} MB)</span>
        `;
        fileList.appendChild(div);
      });
      
      document.getElementById('file-info').classList.remove('hidden');
    }

    document.getElementById('process-btn').addEventListener('click', async () => {
      if (uploadedFiles.length === 0) {
        alert('Please upload a file first');
        return;
      }
      
      if (current_tool.multiple && uploadedFiles.length < 2) {
        alert('Please select at least 2 PDF files to merge');
        return;
      }
      
      const form = new FormData();
      
      if (current_tool.multiple) {
        for (let f of uploadedFiles) {
          form.append('files', f);
        }
      } else {
        form.append('file', uploadedFiles[0]);
      }

      // Get options from the form
      const optionsForm = document.getElementById('options-form');
      const formData = new FormData(optionsForm);
      for (let [key, value] of formData.entries()) {
        form.append(key, value);
      }

      document.getElementById('btn-text').textContent = 'Processing...';
      document.getElementById('spinner').classList.remove('hidden');
      document.getElementById('process-btn').disabled = true;
      document.getElementById('download-link').classList.add('hidden');
      document.getElementById('success-msg').classList.add('hidden');

      try {
        const res = await fetch('/' + current_tool.id, { method: 'POST', body: form });
        const data = await res.json();
        
        if (data.url) {
          document.getElementById('download-link').href = data.url;
          document.getElementById('download-link').classList.remove('hidden');
          document.getElementById('success-msg').classList.remove('hidden');
          document.getElementById('btn-text').textContent = 'Done!';
        } else {
          alert(data.error || 'Processing failed. Please try again.');
          document.getElementById('btn-text').textContent = 'Process PDF';
        }
      } catch (e) {
        alert('Error: ' + e.message);
        document.getElementById('btn-text').textContent = 'Process PDF';
      } finally {
        document.getElementById('spinner').classList.add('hidden');
        document.getElementById('process-btn').disabled = false;
      }
    });
    
    // Close modal on escape key
    document.addEventListener('keydown', (e) => {
      if (e.key === 'Escape') closeModal();
    });
    
    // Close modal on backdrop click
    document.getElementById('tool-modal').addEventListener('click', (e) => {
      if (e.target.id === 'tool-modal') closeModal();
    });
  </script>
</body>
</html>
"""

TOOLS = [
    {"id": "merge", "name": "Merge PDF", "desc": "Combine multiple PDFs into one", "icon": "ph-files", "color": "text-red-600", "bg": "bg-red-100", "multiple": True},
    {"id": "split", "name": "Split PDF", "desc": "Extract pages or split into parts", "icon": "ph-scissors", "color": "text-purple-600", "bg": "bg-purple-100"},
    {"id": "compress", "name": "Compress PDF", "desc": "Reduce file size significantly", "icon": "ph-arrows-in-line-horizontal", "color": "text-green-600", "bg": "bg-green-100"},
    {"id": "word", "name": "PDF to Word", "desc": "Convert to editable .docx", "icon": "ph-microsoft-word-logo", "color": "text-blue-600", "bg": "bg-blue-100"},
    {"id": "excel", "name": "PDF to Excel", "desc": "Extract tables to .xlsx", "icon": "ph-microsoft-excel-logo", "color": "text-green-600", "bg": "bg-green-100"},
    {"id": "ppt", "name": "PDF to PowerPoint", "desc": "Convert to .pptx slides", "icon": "ph-microsoft-powerpoint-logo", "color": "text-orange-600", "bg": "bg-orange-100"},
    {
        "id": "organize", 
        "name": "Organize Pages", 
        "desc": "Reorder, rotate, delete pages", 
        "icon": "ph-squares-four", 
        "color": "text-indigo-600", 
        "bg": "bg-indigo-100",
        "options": '<label class="block"><span class="text-sm font-medium text-gray-700 dark:text-gray-300">Page order (e.g., 3,1,2,5 or leave empty for all pages):</span><input type="text" name="order" class="mt-1 block w-full px-4 py-2 border border-gray-300 dark:border-slate-600 rounded-lg bg-white dark:bg-slate-700 text-gray-900 dark:text-gray-100 focus:ring-2 focus:ring-blue-500" placeholder="1,2,3..."></label>'
    },
]

@app.route('/')
def index():
    response = render_template_string(HTML, tools=TOOLS)
    return response

@app.route('/<tool_id>', methods=['POST'])
def handle_tool(tool_id):
    try:
        if tool_id == "merge":
            files = request.files.getlist('files')
            if len(files) < 2: 
                return jsonify(error="Please select at least 2 PDF files to merge")
            merger = fitz.open()
            for file in files:
                pdf_data = file.read()
                if pdf_data:
                    merger.insert_pdf(fitz.open("pdf", pdf_data))
            output = f"merged_{uuid.uuid4().hex}.pdf"
            merger.save(f"downloads/{output}")
            merger.close()
            return jsonify(url=f"/download/{output}")

        elif tool_id == "compress":
            file = request.files.get('file')
            if not file:
                return jsonify(error="No file uploaded")
            doc = fitz.open("pdf", file.read())
            output = f"compressed_{uuid.uuid4().hex}.pdf"
            doc.save(f"downloads/{output}", garbage=4, deflate=True, clean=True)
            doc.close()
            return jsonify(url=f"/download/{output}")

        elif tool_id == "word":
            file = request.files.get('file')
            if not file:
                return jsonify(error="No file uploaded")
            path = f"uploads/temp_{uuid.uuid4().hex}.pdf"
            file.save(path)
            out = f"word_{uuid.uuid4().hex}.docx"
            cv = Converter(path)
            cv.convert(f"downloads/{out}")
            cv.close()
            os.remove(path)
            return jsonify(url=f"/download/{out}")

        elif tool_id == "excel":
            file = request.files.get('file')
            if not file:
                return jsonify(error="No file uploaded")
            path = f"uploads/temp_{uuid.uuid4().hex}.pdf"
            file.save(path)
            out = f"excel_{uuid.uuid4().hex}.xlsx"
            tables_found = False
            with pdfplumber.open(path) as pdf:
                with pd.ExcelWriter(f"downloads/{out}") as writer:
                    for i, page in enumerate(pdf.pages):
                        tables = page.extract_tables()
                        if tables:
                            tables_found = True
                            for j, table in enumerate(tables):
                                df = pd.DataFrame(table)
                                df.to_excel(writer, sheet_name=f"P{i+1}_T{j+1}", index=False)
                    if not tables_found:
                        # Create empty sheet if no tables found
                        pd.DataFrame({"Note": ["No tables found in this PDF"]}).to_excel(writer, sheet_name="Info", index=False)
            os.remove(path)
            return jsonify(url=f"/download/{out}")

        elif tool_id == "ppt":
            file = request.files.get('file')
            if not file:
                return jsonify(error="No file uploaded")
            path = f"uploads/temp_{uuid.uuid4().hex}.pdf"
            file.save(path)
            prs = Presentation()
            doc = fitz.open(path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = io.BytesIO(pix.tobytes("png"))
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img, Inches(0), Inches(0), width=Inches(10))
            out = f"ppt_{uuid.uuid4().hex}.pptx"
            prs.save(f"downloads/{out}")
            doc.close()
            os.remove(path)
            return jsonify(url=f"/download/{out}")

        elif tool_id == "split":
            file = request.files.get('file')
            if not file:
                return jsonify(error="No file uploaded")
            path = f"uploads/temp_{uuid.uuid4().hex}.pdf"
            file.save(path)
            doc = fitz.open(path)
            zip_path = f"downloads/split_{uuid.uuid4().hex}.zip"
            with zipfile.ZipFile(zip_path, 'w') as zf:
                for i in range(len(doc)):
                    new_doc = fitz.open()
                    new_doc.insert_pdf(doc, from_page=i, to_page=i)
                    buf = io.BytesIO()
                    new_doc.save(buf)
                    zf.writestr(f"page_{i+1}.pdf", buf.getvalue())
                    new_doc.close()
            doc.close()
            os.remove(path)
            return jsonify(url=f"/download/{os.path.basename(zip_path)}")

        elif tool_id == "organize":
            file = request.files.get('file')
            if not file:
                return jsonify(error="No file uploaded")
            order = request.form.get('order', '')
            path = f"uploads/temp_{uuid.uuid4().hex}.pdf"
            file.save(path)
            doc = fitz.open(path)
            
            if order.strip():
                indices = [int(x.strip())-1 for x in order.split(',') if x.strip().isdigit()]
                valid_indices = [i for i in indices if 0 <= i < len(doc)]
                if valid_indices:
                    doc.select(valid_indices)
            
            out = f"organized_{uuid.uuid4().hex}.pdf"
            doc.save(f"downloads/{out}")
            doc.close()
            os.remove(path)
            return jsonify(url=f"/download/{out}")

        else:
            return jsonify(error="Unknown tool")

    except Exception as e:
        return jsonify(error=str(e))

@app.route('/download/<filename>')
def download(filename):
    filepath = f"downloads/{filename}"
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify(error="File not found"), 404

@app.route('/favicon.ico')
def favicon():
    return '', 204

if __name__ == '__main__':
    print("PDFToolz started at http://0.0.0.0:5000")
    app.run(host='0.0.0.0', port=5000, debug=False)
